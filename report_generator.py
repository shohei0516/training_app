# -*- coding: utf-8 -*-
"""
Pitching report generator
"""

from __future__ import annotations

import sys
import subprocess
from pathlib import Path
from typing import Dict, Tuple
from urllib.parse import quote

import pandas as pd
import qrcode
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


# ============================================================
# 1. パス設定
# ============================================================

BASE_DIR = Path(__file__).resolve().parent

EXCEL_FILE = BASE_DIR / "pitching_form_training.xlsx"
PPT_TEMPLATE = BASE_DIR / "ピッチングレポート.pptx"

OUTPUT_DIR = BASE_DIR / "report_output"
OUTPUT_PPTX = OUTPUT_DIR / "pitching_report_auto.pptx"
OUTPUT_PDF = OUTPUT_DIR / "pitching_report_auto.pdf"
QR_FILE = OUTPUT_DIR / "training_qr.png"

TRAINING_MENU_URL = "https://shohei0516.github.io/training_system/html/"

PROBLEM_SCORE_THRESHOLD = 1


# ============================================================
# 2. PPT内の目印名
# ============================================================

TEXT_MARKERS = {
    "comment": "comment_box",
    "qr": "qr_box",
    "total_score": "total_score_box",
}

STAR_MARKERS = {
    "ワインドアップ": "star_windup",
    "アーリー": "star_early",
    "アーリーコッキング": "star_early",
    "レイト": "star_late",
    "レイトコッキング": "star_late",
    "アクセラ": "star_acceleration",
    "アクセラレーション": "star_acceleration",
    "フォロー": "star_follow",
    "フォロースルー": "star_follow",
}

PHASE_FULL_NAME = {
    "ワインドアップ": "ワインドアップ",
    "アーリー": "アーリーコッキング",
    "アーリーコッキング": "アーリーコッキング",
    "レイト": "レイトコッキング",
    "レイトコッキング": "レイトコッキング",
    "アクセラ": "アクセラレーション",
    "アクセラレーション": "アクセラレーション",
    "フォロー": "フォロースルー",
    "フォロースルー": "フォロースルー",
}


# ============================================================
# 3. 基本ユーティリティ
# ============================================================

def normalize_text(value) -> str:
    if pd.isna(value):
        return ""
    return str(value).strip()


def normalize_phase(value) -> str:
    text = normalize_text(value)
    return PHASE_FULL_NAME.get(text, text)


def remove_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def find_shape_by_text(slide, marker_text: str):
    for shape in slide.shapes:
        if hasattr(shape, "text") and normalize_text(shape.text) == marker_text:
            return shape
    return None


def set_text(shape, text: str, font_size: int = 14, bold: bool = False) -> None:
    shape.text = text
    tf = shape.text_frame
    tf.word_wrap = True

    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.name = "Yu Gothic"
            run.font.color.rgb = RGBColor(0, 0, 0)

# ============================================================
# 4. Excel読み込み
# ============================================================

def read_excel_data(excel_path: Path) -> Tuple[pd.DataFrame, pd.DataFrame]:
    if not excel_path.exists():
        raise FileNotFoundError(f"Excelが見つかりません: {excel_path}")

    eval_df = pd.read_excel(excel_path, sheet_name="評価入力")
    link_df = pd.read_excel(excel_path, sheet_name="紐づけ表")

    required_eval_cols = ["チェックID", "フェーズ", "チェック項目", "スコア"]
    required_link_cols = ["チェックID", "トレーニングID", "種目名"]

    for col in required_eval_cols:
        if col not in eval_df.columns:
            raise ValueError(f"評価入力シートに必要な列がありません: {col}")

    for col in required_link_cols:
        if col not in link_df.columns:
            raise ValueError(f"紐づけ表シートに必要な列がありません: {col}")

    eval_df = eval_df.copy()
    eval_df["チェックID"] = eval_df["チェックID"].apply(normalize_text)
    eval_df["スコア"] = pd.to_numeric(eval_df["スコア"], errors="coerce").fillna(0).astype(int)
    eval_df["フェーズ_正式名"] = eval_df["フェーズ"].apply(normalize_phase)

    link_df = link_df.copy()
    link_df["チェックID"] = link_df["チェックID"].apply(normalize_text)
    link_df["トレーニングID"] = link_df["トレーニングID"].apply(normalize_text)

    return eval_df, link_df


# ============================================================
# 5. 推奨トレーニング抽出
# ============================================================

def extract_training(eval_df: pd.DataFrame, link_df: pd.DataFrame) -> pd.DataFrame:
    """
    トレーニング抽出はチェックIDのみで行う。
    スコア1でも2でも、該当チェックIDに紐づくトレーニングをすべて出す。
    """

    problem_ids = (
        eval_df.loc[
            eval_df["スコア"] >= PROBLEM_SCORE_THRESHOLD,
            ["チェックID"]
        ]
        .dropna()
        .drop_duplicates()
    )

    if problem_ids.empty:
        return pd.DataFrame(columns=link_df.columns)

    merged = link_df.merge(
        problem_ids,
        on="チェックID",
        how="inner",
    )

    merged = merged[merged["トレーニングID"].apply(normalize_text) != ""]
    merged = merged.drop_duplicates(subset=["トレーニングID"], keep="first")

    return merged.reset_index(drop=True)


# ============================================================
# 6. コメント自動生成
# ============================================================

def select_interpretation(problems: pd.DataFrame, total_score: int) -> str:
    phases = problems["フェーズ_正式名"].tolist()

    if total_score >= 9:
        return "下肢・体幹・上肢の運動連鎖に乱れがみられ、投球動作全体の効率低下や障害リスク増大が考えられます。"

    if any(p in ["アクセラレーション"] for p in phases):
        return "上肢主導の投球となっている可能性があり、肩・肘への負荷増大やリリースの再現性低下につながる可能性があります。"

    if any(p in ["レイトコッキング"] for p in phases):
        return "ステップ脚の支持性や体幹回旋への連動に課題があり、上肢加速への力の伝達が不十分となっている可能性があります。"

    if any(p in ["ワインドアップ", "アーリーコッキング"] for p in phases):
        return "投球開始からステップ動作にかけての下肢支持性に課題があり、以降の体重移動や回旋動作に影響する可能性があります。"

    return "フォーム全体の再現性に一部課題がみられ、投球動作の安定性低下につながる可能性があります。"


def generate_auto_comment(eval_df: pd.DataFrame, training_df: pd.DataFrame) -> str:
    problems = eval_df[eval_df["スコア"] >= PROBLEM_SCORE_THRESHOLD].copy()

    total_score = int(eval_df["スコア"].sum())
    max_score = int(len(eval_df) * 2)

    if problems.empty:
        return (
            f"総合評価：{total_score} / {max_score}点\n"
            "明らかなフォーム課題はみられません。\n"
            "現在のフォームを維持しつつ、再現性向上を目的とした\n"
            "コンディショニングを継続してください。"
        )

    problems = problems.sort_values(
        ["スコア", "フェーズ_正式名"],
        ascending=[False, True]
    )

    main_phase = (
        problems.groupby("フェーズ_正式名")["スコア"]
        .sum()
        .sort_values(ascending=False)
        .index[0]
    )

    main_items = []
    for _, row in problems.head(4).iterrows():
        main_items.append(f"{row['チェック項目']}（{int(row['スコア'])}点）")

    item_text = "、".join(main_items)
    if len(problems) > 4:
        item_text += " など"

    if total_score <= 4:
        level = "軽度"
    elif total_score <= 8:
        level = "中等度"
    else:
        level = "高度"

    interpretation = select_interpretation(problems, total_score)

    comment = (
        f"総合評価：{total_score} / {max_score}点（{level}のフォーム課題）\n"
        f"主な課題フェーズ：{main_phase}\n"
        f"重要課題：{item_text}\n"
        f"解釈：{interpretation}\n"
        f"推奨：該当項目に対して{len(training_df)}種目のトレーニングを提案します。\n"
        "QRコードから自主トレーニング内容を確認してください。"
    )

    return comment

# ============================================================
# 7. QRコード生成
# ============================================================

def make_training_url(training_df: pd.DataFrame) -> str:
    if training_df.empty:
        return TRAINING_MENU_URL

    ids = [
        normalize_text(x)
        for x in training_df["トレーニングID"].tolist()
        if normalize_text(x)
    ]

    id_text = ",".join(ids)

    return f"{TRAINING_MENU_URL}?ids={quote(id_text)}"


def create_qr(url: str, output_path: Path) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    img = qrcode.make(url)
    img.save(output_path)


# ============================================================
# 8. PPT内テーブルのスコア更新
# ============================================================

def update_score_tables(slide, eval_df: pd.DataFrame) -> None:
    score_by_checkpoint: Dict[str, int] = {}

    for _, row in eval_df.iterrows():
        score_by_checkpoint[normalize_text(row["チェック項目"])] = int(row["スコア"])

    for shape in slide.shapes:
        if not getattr(shape, "has_table", False):
            continue

        table = shape.table

        if len(table.rows) < 2:
            continue

        headers = [
            normalize_text(cell.text).lower()
            for cell in table.rows[0].cells
        ]

        checkpoint_col = None
        score_col = None

        for idx, h in enumerate(headers):
            if h in ["checkpoint", "チェック項目"]:
                checkpoint_col = idx
            if h in ["score", "スコア"]:
                score_col = idx

        if checkpoint_col is None or score_col is None:
            continue

        for row_idx in range(1, len(table.rows)):
            row = table.rows[row_idx]
            checkpoint = normalize_text(row.cells[checkpoint_col].text)

            if checkpoint in score_by_checkpoint:
                row.cells[score_col].text = str(score_by_checkpoint[checkpoint])


# ============================================================
# 9. 赤星追加
# ============================================================

def add_star_at_marker(slide, marker_shape) -> None:
    left = marker_shape.left
    top = marker_shape.top
    width = marker_shape.width
    height = marker_shape.height

    star_box = slide.shapes.add_textbox(left, top, width, height)

    tf = star_box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER

    run = p.add_run()
    run.text = "★"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = RGBColor(220, 0, 0)
    run.font.name = "Yu Gothic"


def add_phase_stars(slide, eval_df: pd.DataFrame) -> None:
    problem_phases = set(
        eval_df.loc[
            eval_df["スコア"] >= PROBLEM_SCORE_THRESHOLD,
            "フェーズ_正式名"
        ].tolist()
    )

    for phase in problem_phases:
        marker_text = STAR_MARKERS.get(phase)

        if not marker_text:
            continue

        marker_shape = find_shape_by_text(slide, marker_text)

        if marker_shape is not None:
            add_star_at_marker(slide, marker_shape)

    for marker_text in set(STAR_MARKERS.values()):
        marker_shape = find_shape_by_text(slide, marker_text)

        if marker_shape is not None:
            remove_shape(marker_shape)


# ============================================================
# 10. PPT作成
# ============================================================

def build_ppt_report(
    ppt_template: Path,
    output_pptx: Path,
    eval_df: pd.DataFrame,
    training_df: pd.DataFrame,
    qr_path: Path,
) -> None:
    if not ppt_template.exists():
        raise FileNotFoundError(f"PPTテンプレートが見つかりません: {ppt_template}")

    prs = Presentation(ppt_template)
    slide = prs.slides[0]

    update_score_tables(slide, eval_df)

    total_score = int(eval_df["スコア"].sum())
    max_score = int(len(eval_df) * 2)
    total_text = f"{total_score} / {max_score}"

    total_shape = find_shape_by_text(slide, TEXT_MARKERS["total_score"])
    if total_shape is not None:
        set_text(total_shape, total_text, font_size=20, bold=True)

    comment = generate_auto_comment(eval_df, training_df)

    comment_shape = find_shape_by_text(slide, TEXT_MARKERS["comment"])
    if comment_shape is not None:
        set_text(comment_shape, comment, font_size=8, bold=False)

    qr_marker = find_shape_by_text(slide, TEXT_MARKERS["qr"])
    if qr_marker is not None:
        left = qr_marker.left
        top = qr_marker.top
        width = qr_marker.width
        height = qr_marker.height

        remove_shape(qr_marker)
        slide.shapes.add_picture(str(qr_path), left, top, width=width, height=height)

    add_phase_stars(slide, eval_df)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_pptx)


# ============================================================
# 11. PDF変換
# ============================================================

def export_pdf_windows_powerpoint(input_pptx: Path, output_pdf: Path) -> bool:
    try:
        import win32com.client  # type: ignore
    except Exception:
        return False

    powerpoint = None
    presentation = None

    try:
        powerpoint = win32com.client.Dispatch("PowerPoint.Application")
        powerpoint.Visible = 1

        presentation = powerpoint.Presentations.Open(str(input_pptx), WithWindow=False)
        presentation.SaveAs(str(output_pdf), 32)
        presentation.Close()
        powerpoint.Quit()

        return True

    except Exception as e:
        print(f"PDF出力に失敗しました: {e}")

        try:
            if presentation:
                presentation.Close()
            if powerpoint:
                powerpoint.Quit()
        except Exception:
            pass

        return False


def export_pdf_libreoffice(input_pptx: Path, output_dir: Path) -> bool:
    try:
        cmd = [
            "soffice",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(input_pptx),
        ]

        result = subprocess.run(cmd, capture_output=True, text=True, check=False)
        return result.returncode == 0

    except Exception:
        return False


# ============================================================
# 12. メイン処理
# ============================================================

def main() -> None:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    eval_df, link_df = read_excel_data(EXCEL_FILE)
    training_df = extract_training(eval_df, link_df)

    training_url = make_training_url(training_df)
    create_qr(training_url, QR_FILE)

    build_ppt_report(
        ppt_template=PPT_TEMPLATE,
        output_pptx=OUTPUT_PPTX,
        eval_df=eval_df,
        training_df=training_df,
        qr_path=QR_FILE,
    )

    print("==== 自動レポート生成 完了 ====")
    print(f"PPTX: {OUTPUT_PPTX}")
    print(f"QR  : {QR_FILE}")
    print(f"URL : {training_url}")
    print(f"推奨トレーニング数: {len(training_df)}")

    pdf_done = export_pdf_windows_powerpoint(OUTPUT_PPTX, OUTPUT_PDF)

    if pdf_done:
        print(f"PDF : {OUTPUT_PDF}")
    else:
        print("PDF出力はスキップされました。PPTXは正常に作成されています。")
        print("PDF化したい場合は、PowerPointでPPTXを開いて『名前を付けて保存 → PDF』で出力できます。")


if __name__ == "__main__":
    try:
        main()
    except Exception as e:
        print("エラーが発生しました。")
        print(e)
        sys.exit(1)